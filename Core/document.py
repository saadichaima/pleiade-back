# Core/document.py
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
from io import BytesIO
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

def _is_pdf(data: bytes) -> bool:
    return data.startswith(b"%PDF")

def extract_text_from_bytes(data: bytes, filename: str = "") -> str:
    name = (filename or "").lower()

    # PDF
    if name.endswith(".pdf") or _is_pdf(data):
        if not data:
            return ""
        doc = fitz.open(stream=data, filetype="pdf")
        return "".join(p.get_text() for p in doc)

    # DOCX
    if name.endswith(".docx"):
        if not data:
            return ""
        d = DocxDocument(BytesIO(data))
        # Iterate all <w:p> elements in the XML (captures paragraphs, tables, SDTs, text boxes…)
        W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        lines = []
        for p_elem in d.element.iter(f"{{{W}}}p"):
            para_text = "".join(t.text or "" for t in p_elem.iter(f"{{{W}}}t"))
            if para_text.strip():
                lines.append(para_text)
        return "\n".join(lines)

    # TXT
    if name.endswith(".txt"):
        return (data or b"").decode("utf-8", errors="ignore")

    # PPTX
    if name.endswith(".pptx"):
        if not data:
            return ""
        prs = Presentation(BytesIO(data))
        out = []
        for s in prs.slides:
            for sh in s.shapes:
                if hasattr(sh, "text"):
                    out.append(sh.text)
        return "\n".join(out)

    # XLSX / XLS
    if name.endswith(".xlsx"):
        return _extract_xlsx_text(data)

    if name.endswith(".xls"):
        return _extract_xls_text(data)

    return ""


def _extract_xlsx_text(data: bytes) -> str:
    """
    Extrait le texte d'un .xlsx avec contexte structurel :
    - Nom de chaque feuille comme en-tête
    - Première ligne utilisée comme en-têtes de colonnes
    - Format : "Colonne: valeur | Colonne2: valeur2"
    - data_only=True : lit les valeurs cachées (pas les formules brutes)
      → si une cellule est None, la valeur n'était pas calculée/cachée dans le fichier
    """
    if not data:
        return ""
    wb = openpyxl.load_workbook(BytesIO(data), data_only=True)
    sections = []
    for sh in wb.worksheets:
        rows = list(sh.iter_rows(values_only=True))
        if not rows:
            continue
        # Première ligne = en-têtes de colonnes
        headers = [
            str(c).strip() if c not in (None, "") else f"Col{i + 1}"
            for i, c in enumerate(rows[0])
        ]
        section_lines = [f"=== Feuille: {sh.title} ==="]
        for row in rows[1:]:
            pairs = []
            for header, cell in zip(headers, row):
                if cell not in (None, ""):
                    val = str(cell).strip()
                    if val:
                        pairs.append(f"{header}: {val}")
            if pairs:
                section_lines.append(" | ".join(pairs))
        if len(section_lines) > 1:
            sections.append("\n".join(section_lines))
    return "\n\n".join(sections)


def _extract_xls_text(data: bytes) -> str:
    """
    Extrait le texte d'un .xls (ancien format Excel) via xlrd.
    Même structure que _extract_xlsx_text.
    """
    if not data:
        return ""
    try:
        import xlrd
    except ImportError:
        return ""
    wb = xlrd.open_workbook(file_contents=data)
    sections = []
    for sh in wb.sheets():
        if sh.nrows == 0:
            continue
        # Première ligne = en-têtes
        headers = [
            str(sh.cell_value(0, c)).strip() or f"Col{c + 1}"
            for c in range(sh.ncols)
        ]
        section_lines = [f"=== Feuille: {sh.name} ==="]
        for r in range(1, sh.nrows):
            pairs = []
            for c in range(sh.ncols):
                cell_type = sh.cell_type(r, c)
                if cell_type == xlrd.XL_CELL_EMPTY:
                    continue
                val = sh.cell_value(r, c)
                # Afficher les entiers sans .0
                if cell_type == xlrd.XL_CELL_NUMBER and val == int(val):
                    str_val = str(int(val))
                else:
                    str_val = str(val).strip()
                if str_val:
                    pairs.append(f"{headers[c]}: {str_val}")
            if pairs:
                section_lines.append(" | ".join(pairs))
        if len(section_lines) > 1:
            sections.append("\n".join(section_lines))
    return "\n\n".join(sections)

def extract_text(file) -> str:
    """Compat rétro: lit le flux entier et appelle extract_text_from_bytes."""
    if not file:
        return ""
    try:
        pos = getattr(file, "tell", lambda: 0)()
        data = file.read() or b""
        # si possible, revenir au début pour ne pas casser l'appelant
        try:
            file.seek(0)
        except Exception:
            pass
        filename = getattr(file, "name", "") or ""
        return extract_text_from_bytes(data, filename)
    except Exception:
        return ""



def chunk_text(text, size=1000, overlap=200):
    words = text.split()
    step = max(1, size - overlap)
    return [" ".join(words[i:i+size]) for i in range(0, len(words), step)]
